import openpyxl
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import math

def _set_shape_alpha(shape, alpha_val):
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        sp_pr = shape._element.spPr
        solid_fill = sp_pr.find(qn('a:solidFill'))
        if solid_fill is not None:
            srgb = solid_fill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_el = srgb.find(qn('a:alpha'))
                if alpha_el is None:
                    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                alpha_el.set('val', str(alpha_val))
    except Exception:
        pass

ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
LOGO_WHITE = os.path.join(ASSETS_DIR, "Royal Purple White Logo.png")
LOGO_SYNTHETIC = os.path.join(ASSETS_DIR, "RPMO_logo_BF_Outline.png")
LOGO_EXPERT_YELLOW = os.path.join(ASSETS_DIR, "RP_Synthetic_Expert_Logo_Yellow_Text.png")
LOGO_EXPERT_BLACK = os.path.join(ASSETS_DIR, "RP_Synthetic_Expert_Logo_Black_Text.png")
LOGO_EXPERT_WHITE = os.path.join(ASSETS_DIR, "rp_synthetic_expert_white.png")
BG_NEVER_SETTLE = os.path.join(ASSETS_DIR, "25-RYP-02147 Employee LinkedIn Thumbnails P1-6.jpg")
IMG_BETTER_OIL = os.path.join(ASSETS_DIR, "Better Oil Starts Here.png")


C = {
    "purple": "4B2D8A",
    "purpleMid": "6B44B8",
    "purpleLight": "9B6FD4",
    "gold": "C8973A",
    "goldLight": "E8B85A",
    "white": "FFFFFF",
    "offWhite": "F8F5FF",
    "lightGray": "F2F2F2",
    "midGray": "94A3B8",
    "darkGray": "334155",
    "dark": "1E1035",
    "green": "22C55E",
    "teal": "0D9488",
}

PRODUCT_MAP = {
    "HMX": "High Mileage",
    "RMS": "High Mileage Syn",
    "RS": "Royal Purple High",
    "RP": "Royal Purple Syn",
    "RSD": "Duralec",
    "11722": "Max-Clean",
    "11755": "Royal Purple Premium",
    "18000": "Max-Atomizer",
}

PRODUCT_FULL_NAMES = {
    "11722": "Max-Clean Fuel System Cleaner",
    "18000": "Max-Atomizer Fuel Injector Cleaner",
    "HMX0W20": "HMX 0W-20 High Mileage",
    "HMX5W20": "HMX 5W-20 High Mileage",
    "HMX5W30": "HMX 5W-30 High Mileage",
    "RMS5W20": "HMX Syn 5W-20 High Mileage",
    "RMS5W30": "HMX Syn 5W-30 High Mileage",
    "RS0W16": "RP High Perf 0W-16",
    "RS0W20": "RP High Perf 0W-20",
    "RS0W40": "RP High Perf 0W-40",
    "RS5W20": "RP High Perf 5W-20",
    "RS5W30": "RP High Perf 5W-30",
    "RS5W40": "RP High Perf 5W-40",
    "RP0W16": "RP Synthetic 0W-16",
    "RP0W20": "RP Synthetic 0W-20",
    "RP0W40": "RP Synthetic 0W-40",
    "RP5W20": "RP Synthetic 5W-20",
    "RP5W30": "RP Synthetic 5W-30",
    "RP5W40": "RP Synthetic 5W-40",
    "RSD15W40": "Duralec Super 15W-40",
    "RSD5W40": "Duralec Super 5W-40",
    "11755": "RP Premium Motor Oil",
}

PRODUCT_DESCRIPTIONS = {
    "High Mileage": "Synthetic motor oil for engines with 75,000+ miles. Reduces oil consumption, revitalizes seals, and removes deposits using Synerlec additive technology.",
    "High Mileage Syn": "Full synthetic high mileage formulation with enhanced seal conditioning and superior wear protection for high-mileage engines.",
    "Royal Purple High": "High-performance full synthetic oil with Synerlec technology for superior wear protection, reduced heat, and increased fuel efficiency.",
    "Royal Purple Syn": "Premium full synthetic oil exceeding API/ILSAC standards. Enhanced film strength minimizes metal-to-metal contact in modern engines.",
    "Duralec": "Premium synthetic diesel engine oil (API CK-4) for emission-controlled engines with DPF, EGR, and SCR systems. Extends drain intervals.",
    "Max-Clean": "High-performance fuel system cleaner and stabilizer. Deeply cleans injectors, carburetors, intake valves, and combustion chambers.",
    "Max-Atomizer": "Advanced fuel injector cleaner for optimized spray patterns and improved combustion efficiency.",
    "Royal Purple Premium": "Premium synthetic motor oil with proprietary Synerlec additive technology for maximum engine protection.",
}

def get_product_display_name(code):
    return PRODUCT_FULL_NAMES.get(code, code)

def get_product_category_desc(category):
    return PRODUCT_DESCRIPTIONS.get(category, "")

def rgb(hex_str):
    return RGBColor.from_string(hex_str)

HEADER_PATTERNS = {
    "date": ["invoice date", "date", "service date", "trans date", "transaction date", "period", "month"],
    "product": ["operation code", "op code", "product", "description", "item", "service", "operation", "sku", "part"],
    "invoices": ["# of invoices", "invoices", "invoice count", "num invoices", "transactions", "oil changes",
                 "ticket count", "tickets", "work orders", "ro count", "repair orders", "# invoices",
                 "number of invoices", "inv count", "total invoices"],
    "revenue": ["total rev", "revenue", "total sales", "net sales", "gross rev", "gross sales",
                "sales rev", "total amount", "net amount", "total $", "gross amount", "sales total",
                "earnings", "income", "proceeds"],
    "avg_rev": ["rev/inv", "avg rev", "average rev", "avg sale", "per invoice", "avg amount",
                "average sale", "rev per", "avg ticket", "average ticket", "per ticket",
                "avg ro", "average order"],
    "vehicles": ["# of vehicles", "vehicles", "vehicle count", "num vehicles", "cars",
                 "unique vehicles", "car count", "unique cars", "vin count"],
    "store": ["store", "location", "shop", "site", "branch", "facility", "installer", "account"],
    "invoice_num": ["invoice #", "invoice num", "invoice number", "inv #", "inv num",
                    "ticket #", "ticket num", "ro #", "ro num", "work order #", "wo #"],
}

SKIP_SHEETS = ["report summary", "summary", "totals", "notes", "instructions", "template", "info",
               "cover", "pivot", "chart", "dashboard", "index",
               "legend", "reference", "lookup", "config", "settings"]


def _safe_float(val, default=0):
    if val is None:
        return default
    try:
        v = str(val).replace("$", "").replace(",", "").strip()
        return float(v)
    except (ValueError, TypeError):
        return default


def _safe_int(val, default=0):
    if val is None:
        return default
    try:
        v = str(val).replace(",", "").strip()
        return int(float(v))
    except (ValueError, TypeError):
        return default


def _find_column_index(header, field):
    patterns = HEADER_PATTERNS.get(field, [])
    header_lower = [str(h).lower().strip() if h else "" for h in header]
    for pattern in patterns:
        for i, h in enumerate(header_lower):
            if pattern == h:
                return i
    for pattern in patterns:
        for i, h in enumerate(header_lower):
            if pattern in h:
                return i
    return None


def _find_header_row(rows, max_scan=10):
    keywords = ["invoice", "date", "revenue", "product", "operation", "sales",
                "amount", "total", "store", "location", "description", "qty",
                "vehicles", "transactions", "tickets", "sku", "service"]
    best_idx = 0
    best_score = 0
    for i, row in enumerate(rows[:max_scan]):
        if not row:
            continue
        row_strs = [str(c).lower().strip() if c else "" for c in row]
        score = sum(1 for s in row_strs if any(kw in s for kw in keywords))
        if score > best_score:
            best_score = score
            best_idx = i
    return best_idx if best_score >= 2 else 0


def _parse_single_date(date_val):
    if date_val is None:
        return None
    if hasattr(date_val, 'strftime'):
        return date_val
    if isinstance(date_val, str):
        date_val = date_val.strip()
        if not date_val:
            return None
        import re
        from datetime import datetime as dt
        m = re.match(r'(\d{1,2})[/\-](\d{1,2})[/\-](\d{2,4})', date_val)
        if m:
            for fmt in ["%m/%d/%Y", "%m-%d-%Y", "%m/%d/%y", "%m-%d-%y",
                        "%d/%m/%Y", "%d-%m-%Y", "%Y-%m-%d"]:
                try:
                    return dt.strptime(date_val, fmt)
                except ValueError:
                    continue
        parts = date_val.split()
        if len(parts) >= 3:
            for fmt in ["%B %d, %Y", "%B %d %Y"]:
                try:
                    return dt.strptime(date_val.replace(",", "").strip(), "%B %d %Y")
                except ValueError:
                    pass
    return None


def _detect_date_from_rows(data_rows, col_map):
    date_idx = col_map.get("date")
    if date_idx is None:
        return None
    all_dates = []
    for row in data_rows:
        if date_idx >= len(row):
            continue
        d = _parse_single_date(row[date_idx])
        if d:
            all_dates.append(d)
    if not all_dates:
        return None
    min_d = min(all_dates)
    max_d = max(all_dates)
    if min_d.year == max_d.year and min_d.month == max_d.month:
        return min_d.strftime("%B %Y")
    elif min_d.year == max_d.year:
        return str(min_d.year)
    else:
        return f"{min_d.strftime('%B %Y')} – {max_d.strftime('%B %Y')}"


def _detect_date_from_sheet(ws_title, rows):
    import re
    for text in [ws_title] + [str(c) for r in rows[:5] if r for c in r if c]:
        months = ["january", "february", "march", "april", "may", "june",
                  "july", "august", "september", "october", "november", "december"]
        text_lower = text.lower()
        for month in months:
            if month in text_lower:
                m = re.search(r'(\d{4})', text)
                if m:
                    return f"{month.capitalize()} {m.group(1)}"
    return None


def _get_val(row, idx, default=None):
    if idx is None or idx >= len(row):
        return default
    return row[idx]


def _is_numeric_column(rows, col_idx, sample_size=10):
    count = 0
    numeric = 0
    for row in rows[:sample_size]:
        if col_idx >= len(row):
            continue
        val = row[col_idx]
        if val is None:
            continue
        count += 1
        try:
            float(str(val).replace("$", "").replace(",", ""))
            numeric += 1
        except (ValueError, TypeError):
            pass
    return numeric > 0 and numeric / max(count, 1) >= 0.5


def _detect_revenue_column(header, data_rows):
    header_lower = [str(h).lower().strip() if h else "" for h in header]
    candidates = []
    for i, h in enumerate(header_lower):
        if not h:
            continue
        if any(kw in h for kw in ["rev", "sales", "amount", "total", "earnings", "income", "gross", "net"]):
            if any(skip in h for skip in ["date", "name", "desc", "code", "model", "make", "year", "vin"]):
                continue
            if _is_numeric_column(data_rows, i):
                total = sum(_safe_float(_get_val(r, i)) for r in data_rows[:50])
                candidates.append((i, total, h))
    if candidates:
        candidates.sort(key=lambda x: -x[1])
        return candidates[0][0]
    for i in range(len(header)):
        if _is_numeric_column(data_rows, i):
            total = sum(_safe_float(_get_val(r, i)) for r in data_rows[:50])
            if total > 100:
                return i
    return None


MC_CODE = "11722"
RP_OIL_PREFIXES = ("RP", "RS", "HMX", "RMS", "RSD")


def _extract_product_code(row, product_idx):
    if product_idx is None:
        return ""
    raw = _get_val(row, product_idx)
    if not raw:
        return ""
    op_desc = str(raw).strip()
    return op_desc.split(" - ")[0].strip() if " - " in op_desc else op_desc.strip()


def _is_rp_oil_code(code):
    c = code.upper()
    return any(c.startswith(p.upper()) for p in RP_OIL_PREFIXES) and c != MC_CODE and c != "18000"


def _group_invoices(data_rows, col_map):
    from collections import defaultdict
    inv_num_idx = col_map.get("invoice_num")
    product_idx = col_map.get("product")
    revenue_idx = col_map.get("revenue")
    date_idx = col_map.get("date")
    vehicle_idx = col_map.get("vehicles")

    groups = defaultdict(lambda: {"rows": [], "codes": [], "revenue": 0, "vehicles": 0})

    for row in data_rows:
        if inv_num_idx is not None:
            inv_key = _get_val(row, inv_num_idx)
            if inv_key is None or str(inv_key).strip() == "":
                date_val = str(_get_val(row, date_idx, "")) if date_idx is not None else ""
                rev_val = _safe_float(_get_val(row, revenue_idx))
                veh_val = str(_get_val(row, vehicle_idx, "")) if vehicle_idx is not None else ""
                inv_key = ("_fallback_", date_val, rev_val, veh_val)
        else:
            date_val = str(_get_val(row, date_idx, "")) if date_idx is not None else ""
            rev_val = _safe_float(_get_val(row, revenue_idx))
            veh_val = str(_get_val(row, vehicle_idx, "")) if vehicle_idx is not None else ""
            inv_key = (date_val, rev_val, veh_val)

        code = _extract_product_code(row, product_idx)
        rev = _safe_float(_get_val(row, revenue_idx))
        veh = _safe_int(_get_val(row, vehicle_idx))

        groups[inv_key]["rows"].append(row)
        if code:
            groups[inv_key]["codes"].append(code)
        groups[inv_key]["revenue"] = rev
        groups[inv_key]["vehicles"] = max(groups[inv_key]["vehicles"], veh)

    return dict(groups)


def _parse_single_store_sheet(sheet_name, rows):
    header_row_idx = _find_header_row(rows)
    header = rows[header_row_idx]

    col_map = {}
    for field in ["date", "product", "invoices", "revenue", "avg_rev", "vehicles", "invoice_num"]:
        idx = _find_column_index(header, field)
        if idx is not None:
            col_map[field] = idx

    if "revenue" not in col_map:
        rev_idx = _detect_revenue_column(header, rows[header_row_idx + 1:])
        if rev_idx is not None:
            col_map["revenue"] = rev_idx

    if "revenue" not in col_map and "invoices" not in col_map:
        return None

    all_data_rows = rows[header_row_idx + 1:]
    first_col = col_map.get("date", col_map.get("product", 0))

    data_rows = []
    for r in all_data_rows:
        if len(r) <= first_col:
            continue
        if r[first_col] is not None:
            data_rows.append(r)

    if not data_rows:
        non_empty = [r for r in all_data_rows if any(c is not None for c in r)]
        if non_empty:
            data_rows = non_empty

    if not data_rows:
        return None

    last_row = all_data_rows[-1] if all_data_rows else None
    totals_row = None
    if last_row:
        first_empty = (len(last_row) <= first_col or last_row[first_col] is None)
        has_numbers = any(_safe_float(_get_val(last_row, col_map.get(f))) > 0
                         for f in ["revenue", "invoices"] if f in col_map)
        if first_empty and has_numbers:
            totals_row = last_row
            if last_row in data_rows:
                data_rows.remove(last_row)

    invoice_groups = _group_invoices(data_rows, col_map)

    product_idx = col_map.get("product")
    sorted_prefixes = sorted(PRODUCT_MAP.keys(), key=len, reverse=True)

    product_line_count = {}
    dedup_revenue = 0
    dedup_invoices = len(invoice_groups)
    dedup_vehicles = 0

    mc_total = 0
    mc_with_rp_oil = 0
    mc_solo_in_data = 0
    mc_revenue_total = 0
    non_mc_revenue_total = 0
    mc_invoice_revenue = 0
    non_mc_invoice_revenue = 0

    for inv_key, inv_data in invoice_groups.items():
        inv_rev = inv_data["revenue"]
        dedup_revenue += inv_rev
        dedup_vehicles += inv_data["vehicles"]

        codes = inv_data["codes"]
        for code in set(codes):
            product_line_count[code] = product_line_count.get(code, 0) + 1

        has_mc = MC_CODE in codes
        has_rp_oil = any(_is_rp_oil_code(c) for c in codes)

        if has_mc:
            mc_total += 1
            mc_invoice_revenue += inv_rev
            if has_rp_oil:
                mc_with_rp_oil += 1
            elif len(codes) <= 1:
                mc_solo_in_data += 1
        else:
            non_mc_invoice_revenue += inv_rev

    mc_non_rp = mc_total - mc_with_rp_oil
    mc_avg_ticket = mc_invoice_revenue / mc_total if mc_total else 0
    non_mc_count = dedup_invoices - mc_total
    non_mc_avg_ticket = non_mc_invoice_revenue / non_mc_count if non_mc_count else 0
    mc_ticket_lift = mc_avg_ticket - non_mc_avg_ticket

    product_breakdown = []
    for code in sorted(product_line_count.keys(), key=lambda c: -product_line_count.get(c, 0)):
        cat = "Other"
        for prefix in sorted_prefixes:
            if code.upper().startswith(prefix.upper()):
                cat = PRODUCT_MAP[prefix]
                break
        raw_rev = sum(_safe_float(_get_val(r, col_map.get("revenue")))
                      for r in data_rows if _extract_product_code(r, product_idx) == code)
        product_breakdown.append({
            "code": code,
            "category": cat,
            "revenue": round(raw_rev, 2),
            "lineCount": product_line_count[code],
        })

    top_product = product_breakdown[0]["category"] if product_breakdown else "N/A"

    avg_rev_inv = dedup_revenue / dedup_invoices if dedup_invoices else 0

    mc_attachment_rate = mc_total / dedup_invoices * 100 if dedup_invoices else 0

    return {
        "name": sheet_name,
        "invoices": int(dedup_invoices),
        "vehicles": int(dedup_vehicles),
        "totalRevenue": round(float(dedup_revenue), 2),
        "avgRevPerInvoice": round(float(avg_rev_inv), 2),
        "topProduct": top_product,
        "productBreakdown": product_breakdown,
        "rawLineCount": len(data_rows),
        "maxClean": {
            "total": mc_total,
            "withRpOil": mc_with_rp_oil,
            "withNonRpOil": mc_non_rp,
            "soloInData": mc_solo_in_data,
            "attachmentRate": round(mc_attachment_rate, 1),
            "avgTicket": round(mc_avg_ticket, 2),
            "nonMcAvgTicket": round(non_mc_avg_ticket, 2),
            "ticketLift": round(mc_ticket_lift, 2),
        },
        "_col_map": col_map,
        "_date_rows": data_rows,
    }


def _parse_consolidated_sheet(sheet_name, rows):
    header_row_idx = _find_header_row(rows)
    header = rows[header_row_idx]

    store_idx = _find_column_index(header, "store")
    if store_idx is None:
        return []

    col_map = {}
    for field in ["date", "product", "invoices", "revenue", "avg_rev", "vehicles", "store", "invoice_num"]:
        idx = _find_column_index(header, field)
        if idx is not None:
            col_map[field] = idx

    if "revenue" not in col_map:
        rev_idx = _detect_revenue_column(header, rows[header_row_idx + 1:])
        if rev_idx is not None:
            col_map["revenue"] = rev_idx

    if "revenue" not in col_map:
        return []

    all_data_rows = rows[header_row_idx + 1:]
    data_rows = [r for r in all_data_rows
                 if len(r) > store_idx and r[store_idx] is not None
                 and str(r[store_idx]).strip()]

    store_rows = {}
    for row in data_rows:
        store_name = str(_get_val(row, store_idx, "")).strip()
        if not store_name or store_name.lower() in ["total", "totals", "grand total", "sum"]:
            continue
        store_rows.setdefault(store_name, []).append(row)

    sorted_prefixes = sorted(PRODUCT_MAP.keys(), key=len, reverse=True)
    product_idx = col_map.get("product")
    stores = []

    for sname, s_rows in store_rows.items():
        invoice_groups = _group_invoices(s_rows, col_map)

        product_line_count = {}
        dedup_revenue = 0
        dedup_invoices = len(invoice_groups)
        dedup_vehicles = 0

        mc_total = 0
        mc_with_rp_oil = 0
        mc_solo_in_data = 0
        mc_invoice_revenue = 0
        non_mc_invoice_revenue = 0

        for inv_key, inv_data in invoice_groups.items():
            inv_rev = inv_data["revenue"]
            dedup_revenue += inv_rev
            dedup_vehicles += inv_data["vehicles"]

            codes = inv_data["codes"]
            for code in set(codes):
                product_line_count[code] = product_line_count.get(code, 0) + 1

            has_mc = MC_CODE in codes
            has_rp_oil = any(_is_rp_oil_code(c) for c in codes)

            if has_mc:
                mc_total += 1
                mc_invoice_revenue += inv_rev
                if has_rp_oil:
                    mc_with_rp_oil += 1
                elif len(codes) <= 1:
                    mc_solo_in_data += 1
            else:
                non_mc_invoice_revenue += inv_rev

        mc_non_rp = mc_total - mc_with_rp_oil
        mc_avg_ticket = mc_invoice_revenue / mc_total if mc_total else 0
        non_mc_count = dedup_invoices - mc_total
        non_mc_avg_ticket = non_mc_invoice_revenue / non_mc_count if non_mc_count else 0
        mc_ticket_lift = mc_avg_ticket - non_mc_avg_ticket
        mc_attachment_rate = mc_total / dedup_invoices * 100 if dedup_invoices else 0

        product_breakdown = []
        for code in sorted(product_line_count.keys(), key=lambda c: -product_line_count.get(c, 0)):
            cat = "Other"
            for prefix in sorted_prefixes:
                if code.upper().startswith(prefix.upper()):
                    cat = PRODUCT_MAP[prefix]
                    break
            raw_rev = sum(_safe_float(_get_val(r, col_map.get("revenue")))
                          for r in s_rows if _extract_product_code(r, product_idx) == code)
            product_breakdown.append({
                "code": code,
                "category": cat,
                "revenue": round(raw_rev, 2),
                "lineCount": product_line_count[code],
            })

        top_product = product_breakdown[0]["category"] if product_breakdown else "N/A"
        avg_rev_inv = dedup_revenue / dedup_invoices if dedup_invoices else 0

        stores.append({
            "name": sname,
            "invoices": int(dedup_invoices),
            "vehicles": int(dedup_vehicles),
            "totalRevenue": round(float(dedup_revenue), 2),
            "avgRevPerInvoice": round(float(avg_rev_inv), 2),
            "topProduct": top_product,
            "productBreakdown": product_breakdown,
            "rawLineCount": len(s_rows),
            "maxClean": {
                "total": mc_total,
                "withRpOil": mc_with_rp_oil,
                "withNonRpOil": mc_non_rp,
                "soloInData": mc_solo_in_data,
                "attachmentRate": round(mc_attachment_rate, 1),
                "avgTicket": round(mc_avg_ticket, 2),
                "nonMcAvgTicket": round(non_mc_avg_ticket, 2),
                "ticketLift": round(mc_ticket_lift, 2),
            },
        })

    return stores


def parse_excel(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    stores = []
    month_year = None

    skip_lower = set(SKIP_SHEETS)
    data_sheets = [s for s in wb.sheetnames if s.lower().strip() not in skip_lower]

    if len(data_sheets) == 1:
        ws = wb[data_sheets[0]]
        rows = list(ws.iter_rows(values_only=True))
        if len(rows) >= 2:
            header_idx = _find_header_row(rows)
            header = rows[header_idx]
            store_idx = _find_column_index(header, "store")

            if store_idx is not None:
                consolidated = _parse_consolidated_sheet(data_sheets[0], rows)
                if consolidated:
                    stores.extend(consolidated)
                    month_year = _detect_date_from_rows(rows[header_idx + 1:],
                                                        {"date": _find_column_index(header, "date")})
                    if not month_year:
                        month_year = _detect_date_from_sheet(data_sheets[0], rows)
            else:
                result = _parse_single_store_sheet(data_sheets[0], rows)
                if result:
                    month_year = _detect_date_from_rows(
                        result.pop("_date_rows", []),
                        result.pop("_col_map", {}))
                    stores.append(result)
    else:
        for sheet_name in data_sheets:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if len(rows) < 2:
                continue

            header_idx = _find_header_row(rows)
            header = rows[header_idx]
            store_idx = _find_column_index(header, "store")

            if store_idx is not None:
                consolidated = _parse_consolidated_sheet(sheet_name, rows)
                if consolidated:
                    stores.extend(consolidated)
                    if not month_year:
                        date_idx = _find_column_index(header, "date")
                        month_year = _detect_date_from_rows(rows[header_idx + 1:],
                                                            {"date": date_idx})
                    continue

            result = _parse_single_store_sheet(sheet_name, rows)
            if result:
                if not month_year:
                    month_year = _detect_date_from_rows(
                        result.get("_date_rows", []),
                        result.get("_col_map", {}))
                    if not month_year:
                        month_year = _detect_date_from_sheet(sheet_name, rows)
                result.pop("_date_rows", None)
                result.pop("_col_map", None)
                stores.append(result)

    stores.sort(key=lambda s: -s["totalRevenue"])
    for i, s in enumerate(stores):
        s["rank"] = i + 1

    if not month_year:
        from datetime import datetime
        month_year = datetime.now().strftime("%B %Y")

    if not stores:
        raise ValueError(
            "No store data found in the Excel file. The app looks for columns like "
            "'Revenue', 'Total Rev', 'Sales', 'Invoices', 'Product', etc. "
            "Make sure your report has recognizable column headers."
        )

    return stores, month_year


def add_slide_background(slide, color=C["offWhite"]):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_top_bar(slide):
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(C["purple"])
    shape.line.fill.background()


def add_footer(slide, page_num, total_slides):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(5.33), Inches(10), Inches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(C["purple"])
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"  {page_num} / {total_slides}"
    run.font.size = Pt(7)
    run.font.color.rgb = rgb(C["white"])
    run.font.name = "Calibri"


def add_royal_purple_badge(slide):
    txBox = slide.shapes.add_textbox(
        Inches(7.8), Inches(0.12), Inches(2.0), Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "ROYAL PURPLE"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = rgb(C["gold"])
    run.font.name = "Calibri"


def add_slide_header(slide, title, subtitle=None):
    add_top_bar(slide)
    add_royal_purple_badge(slide)

    accent = slide.shapes.add_shape(
        1, Inches(0.4), Inches(0.7), Inches(0.07), Inches(0.45)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(C["gold"])
    accent.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.65), Inches(8.0), Inches(0.5)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = rgb(C["purple"])
    run.font.name = "Calibri"

    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.3)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(10)
        run2.font.color.rgb = rgb(C["midGray"])
        run2.font.name = "Calibri"


def _add_metric(slide, x, y, w, value, label, sub_label=None):
    val_str = str(value)
    font_size = 20 if len(val_str) > 8 else 26
    vb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.45))
    tf = vb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = val_str
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = rgb(C["purple"])
    r.font.name = "Calibri"

    lb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.42), Inches(w), Inches(0.22))
    tf2 = lb.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(9)
    r2.font.color.rgb = rgb(C["darkGray"])
    r2.font.name = "Calibri"

    if sub_label:
        sb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.62), Inches(w), Inches(0.18))
        tf3 = sb.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = sub_label
        r3.font.size = Pt(7)
        r3.font.color.rgb = rgb(C["midGray"])
        r3.font.name = "Calibri"


def _add_thin_divider(slide, x, y, w, color=C["lightGray"]):
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(color)
    line.line.fill.background()


def fmt_currency(val):
    if val >= 1_000_000:
        return f"${val/1_000_000:.2f}M"
    elif val >= 1_000:
        return f"${val:,.0f}"
    else:
        return f"${val:.2f}"


def fmt_number(val):
    return f"{val:,}"


def build_cover_slide(prs, stores, month_year, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, C["dark"])

    if os.path.exists(BG_NEVER_SETTLE):
        slide.shapes.add_picture(
            BG_NEVER_SETTLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
        )

    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = rgb(C["dark"])
    _set_shape_alpha(overlay, 60000)
    overlay.line.fill.background()

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(C["gold"])
    bar.line.fill.background()

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Partnership Hub Report"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = rgb(C["white"])
    run2.font.name = "Calibri"

    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(6), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"{month_year} | {len(stores)} Locations"
    run3.font.size = Pt(14)
    run3.font.color.rgb = rgb(C["purpleLight"])
    run3.font.name = "Calibri"

    ns_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(3), Inches(0.35))
    tf_ns = ns_box.text_frame
    p_ns = tf_ns.paragraphs[0]
    r_ns = p_ns.add_run()
    r_ns.text = "NEVER SETTLE"
    r_ns.font.size = Pt(16)
    r_ns.font.bold = True
    r_ns.font.color.rgb = rgb(C["white"])
    r_ns.font.name = "Calibri"

    total_rev = sum(s["totalRevenue"] for s in stores)
    total_inv = sum(s["invoices"] for s in stores)
    avg_rev = total_rev / total_inv if total_inv else 0

    stat_x = 6.8
    stat_y = 1.2
    stat_w = 2.8
    stat_h = 2.2

    stat_bg = slide.shapes.add_shape(1, Inches(stat_x), Inches(stat_y), Inches(stat_w), Inches(stat_h))
    stat_bg.fill.solid()
    stat_bg.fill.fore_color.rgb = rgb(C["purple"])
    stat_bg.line.fill.background()

    gold_top = slide.shapes.add_shape(1, Inches(stat_x), Inches(stat_y), Inches(stat_w), Inches(0.06))
    gold_top.fill.solid()
    gold_top.fill.fore_color.rgb = rgb(C["gold"])
    gold_top.line.fill.background()

    lbl = slide.shapes.add_textbox(Inches(stat_x + 0.15), Inches(stat_y + 0.15), Inches(stat_w - 0.3), Inches(0.25))
    tf_l = lbl.text_frame
    p_l = tf_l.paragraphs[0]
    r_l = p_l.add_run()
    r_l.text = f"{month_year} Summary"
    r_l.font.size = Pt(10)
    r_l.font.bold = True
    r_l.font.color.rgb = rgb(C["goldLight"])
    r_l.font.name = "Calibri"

    items = [
        (fmt_currency(total_rev), "Total Revenue"),
        (fmt_number(total_inv), "Oil Changes"),
        (f"${avg_rev:.2f}", "Avg Rev/Invoice"),
        (stores[0]["name"] if stores else "N/A", "Top Store"),
    ]
    for i, (val, lab) in enumerate(items):
        iy = stat_y + 0.5 + i * 0.42
        vb = slide.shapes.add_textbox(Inches(stat_x + 0.2), Inches(iy), Inches(stat_w - 0.4), Inches(0.22))
        tf_v = vb.text_frame
        p_v = tf_v.paragraphs[0]
        r_v = p_v.add_run()
        r_v.text = val
        r_v.font.size = Pt(14)
        r_v.font.bold = True
        r_v.font.color.rgb = rgb(C["white"])
        r_v.font.name = "Calibri"

        lb = slide.shapes.add_textbox(Inches(stat_x + 0.2), Inches(iy + 0.2), Inches(stat_w - 0.4), Inches(0.15))
        tf_lb = lb.text_frame
        p_lb = tf_lb.paragraphs[0]
        r_lb = p_lb.add_run()
        r_lb.text = lab
        r_lb.font.size = Pt(7)
        r_lb.font.color.rgb = rgb(C["purpleLight"])
        r_lb.font.name = "Calibri"

    from datetime import datetime as _dt
    prepared_date = _dt.now().strftime("%B %d, %Y")

    footer_txt = slide.shapes.add_textbox(Inches(0.8), Inches(4.55), Inches(8), Inches(0.3))
    tf_f = footer_txt.text_frame
    p_f = tf_f.paragraphs[0]
    r_f = p_f.add_run()
    r_f.text = f"Prepared {prepared_date}"
    r_f.font.size = Pt(9)
    r_f.font.color.rgb = rgb(C["purpleLight"])
    r_f.font.name = "Calibri"

    footer_txt2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.85), Inches(8), Inches(0.3))
    tf_f2 = footer_txt2.text_frame
    p_f2 = tf_f2.paragraphs[0]
    r_f2 = p_f2.add_run()
    r_f2.text = "ThrottlePro — More Cars. More Loyalty. Less Stress."
    r_f2.font.size = Pt(8)
    r_f2.font.color.rgb = rgb(C["midGray"])
    r_f2.font.name = "Calibri"

    add_footer(slide, 1, total_slides)


def build_toc_slide(prs, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_slide_header(slide, "Table of Contents", "Report sections overview")
    add_footer(slide, 2, total_slides)

    sections = [
        ("1", "Executive Summary", "KPIs & observations"),
        ("2", "Revenue Overview", "Total revenue breakdown"),
        ("3", "Store Rankings", "Performance table & matrix"),
        ("4", "Product Mix", "Category analysis"),
        ("5", "Product Deep Dives", "Per-category performance"),
        ("6", "Store Deep Dives", "Individual store detail"),
        ("7", "Next Steps", "Recommendations"),
    ]
    for i, (num, title, desc) in enumerate(sections):
        y = 1.5 + i * 0.5
        _add_thin_divider(slide, 0.5, y, 9.0)

        nb = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.06), Inches(0.5), Inches(0.35))
        tf = nb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = num
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = rgb(C["gold"])
        r.font.name = "Calibri"

        tb = slide.shapes.add_textbox(Inches(1.1), Inches(y + 0.08), Inches(3.5), Inches(0.25))
        tf2 = tb.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = title
        r2.font.size = Pt(12)
        r2.font.bold = True
        r2.font.color.rgb = rgb(C["purple"])
        r2.font.name = "Calibri"

        db = slide.shapes.add_textbox(Inches(5.0), Inches(y + 0.1), Inches(4.5), Inches(0.25))
        tf3 = db.text_frame
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = desc
        r3.font.size = Pt(10)
        r3.font.color.rgb = rgb(C["midGray"])
        r3.font.name = "Calibri"


def build_exec_summary_kpis(prs, stores, month_year, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_slide_header(slide, "Executive Summary", f"Key performance indicators — {month_year}")
    add_footer(slide, 3, total_slides)

    total_rev = sum(s["totalRevenue"] for s in stores)
    total_inv = sum(s["invoices"] for s in stores)
    avg_rev = total_rev / total_inv if total_inv else 0
    total_veh = sum(s["vehicles"] for s in stores)

    metrics = [
        (fmt_currency(total_rev), "Total Revenue", f"Across {len(stores)} locations"),
        (fmt_number(total_inv), "Total Oil Changes", f"{month_year}"),
        (f"${avg_rev:.2f}", "Avg Rev / Invoice", "Network average"),
        (fmt_number(total_veh), "Unique Vehicles", f"{month_year}"),
    ]
    mw = 2.1
    gap = 0.2
    sx = (10 - (4 * mw + 3 * gap)) / 2
    for i, (val, lbl, sub) in enumerate(metrics):
        _add_metric(slide, sx + i * (mw + gap), 1.55, mw, val, lbl, sub)

    _add_thin_divider(slide, 0.5, 2.55, 9.0)

    top8 = stores[:8]
    max_rev = top8[0]["totalRevenue"] if top8 else 1

    lbl_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4), Inches(0.3))
    tf = lbl_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Revenue Leaderboard"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = rgb(C["purple"])
    r.font.name = "Calibri"

    for i, store in enumerate(top8):
        y = 3.05 + i * 0.3
        bar_w = max(0.3, (store["totalRevenue"] / max_rev) * 5.0)

        nm = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.8), Inches(0.26))
        tf_n = nm.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.RIGHT
        r_n = p_n.add_run()
        r_n.text = store["name"]
        r_n.font.size = Pt(8)
        r_n.font.color.rgb = rgb(C["darkGray"])
        r_n.font.name = "Calibri"

        bar = slide.shapes.add_shape(1, Inches(3.4), Inches(y + 0.04), Inches(bar_w), Inches(0.18))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(C["purple"])
        bar.line.fill.background()

        vb = slide.shapes.add_textbox(Inches(3.4 + bar_w + 0.08), Inches(y), Inches(1.5), Inches(0.26))
        tf_v = vb.text_frame
        p_v = tf_v.paragraphs[0]
        r_v = p_v.add_run()
        r_v.text = fmt_currency(store["totalRevenue"])
        r_v.font.size = Pt(8)
        r_v.font.bold = True
        r_v.font.color.rgb = rgb(C["purple"])
        r_v.font.name = "Calibri"


def build_exec_observations(prs, stores, month_year, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_slide_header(slide, "Executive Observations", f"Key insights — {month_year}")
    add_footer(slide, 4, total_slides)

    total_rev = sum(s["totalRevenue"] for s in stores)
    total_inv = sum(s["invoices"] for s in stores)
    avg_rev = total_rev / total_inv if total_inv else 0
    top_store = stores[0] if stores else None
    bottom_store = stores[-1] if stores else None

    hm_rev = sum(
        pb["revenue"] for s in stores for pb in s["productBreakdown"]
        if pb["category"] in ("High Mileage", "High Mileage Syn")
    )
    hm_pct = (hm_rev / total_rev * 100) if total_rev else 0

    observations = [
        (
            "Network Revenue",
            f"Total network revenue of {fmt_currency(total_rev)} across {len(stores)} locations with {fmt_number(total_inv)} oil changes at an average of ${avg_rev:.2f} per invoice.",
        ),
        (
            "Top Performer",
            f"{top_store['name']} leads with {fmt_currency(top_store['totalRevenue'])} in revenue ({top_store['totalRevenue']/total_rev*100:.1f}% of network total) and {fmt_number(top_store['invoices'])} oil changes." if top_store else "N/A",
        ),
        (
            "Product Mix Insight",
            f"High Mileage products account for {hm_pct:.1f}% of total revenue, indicating strong upsell penetration across the network.",
        ),
        (
            "Growth Opportunity",
            f"{bottom_store['name']} has the lowest revenue at {fmt_currency(bottom_store['totalRevenue'])} — targeted promotions could lift volume." if bottom_store else "N/A",
        ),
    ]

    for i, (title, body) in enumerate(observations):
        y = 1.55 + i * 0.92
        _add_thin_divider(slide, 0.5, y, 9.0)

        nb = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.08), Inches(0.35), Inches(0.3))
        tf_n = nb.text_frame
        p_n = tf_n.paragraphs[0]
        r_n = p_n.add_run()
        r_n.text = str(i + 1)
        r_n.font.size = Pt(16)
        r_n.font.bold = True
        r_n.font.color.rgb = rgb(C["gold"])
        r_n.font.name = "Calibri"

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.08), Inches(8.5), Inches(0.25))
        tf_t = tb.text_frame
        p_t = tf_t.paragraphs[0]
        r_t = p_t.add_run()
        r_t.text = title
        r_t.font.size = Pt(12)
        r_t.font.bold = True
        r_t.font.color.rgb = rgb(C["purple"])
        r_t.font.name = "Calibri"

        bb = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.36), Inches(8.5), Inches(0.5))
        tf_b = bb.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        r_b = p_b.add_run()
        r_b.text = body
        r_b.font.size = Pt(9)
        r_b.font.color.rgb = rgb(C["darkGray"])
        r_b.font.name = "Calibri"


def build_revenue_overview(prs, stores, month_year, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_slide_header(slide, "Revenue Overview", f"Total revenue distribution — {month_year}")
    add_footer(slide, 5, total_slides)

    total_rev = sum(s["totalRevenue"] for s in stores)
    total_inv = sum(s["invoices"] for s in stores)
    avg_rev = total_rev / total_inv if total_inv else 0

    _add_metric(slide, 0.5, 1.55, 2.8, fmt_currency(total_rev), "Total Network Revenue", f"{len(stores)} Locations | {month_year}")
    _add_metric(slide, 3.5, 1.55, 2.0, fmt_number(total_inv), "Oil Changes", "")
    _add_metric(slide, 5.7, 1.55, 2.0, f"${avg_rev:.2f}", "Avg Rev/Invoice", "")
    _add_metric(slide, 7.8, 1.55, 1.8, fmt_number(sum(s["vehicles"] for s in stores)), "Vehicles", "")

    _add_thin_divider(slide, 0.5, 2.55, 9.0)

    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(5), Inches(0.3))
    tf_h = hdr.text_frame
    p_h = tf_h.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = "Revenue by Store"
    r_h.font.size = Pt(11)
    r_h.font.bold = True
    r_h.font.color.rgb = rgb(C["purple"])
    r_h.font.name = "Calibri"

    max_rows = int((5.33 - 3.05) / 0.28)
    visible_stores = stores[:max_rows]
    for i, store in enumerate(visible_stores):
        sy = 3.05 + i * 0.28
        pct = store["totalRevenue"] / total_rev * 100 if total_rev else 0

        nm = slide.shapes.add_textbox(Inches(0.5), Inches(sy), Inches(2.8), Inches(0.25))
        tf_n = nm.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.RIGHT
        r_n = p_n.add_run()
        r_n.text = store["name"]
        r_n.font.size = Pt(8)
        r_n.font.color.rgb = rgb(C["darkGray"])
        r_n.font.name = "Calibri"

        bar_w = max(0.15, pct / 100 * 4.0)
        bar = slide.shapes.add_shape(1, Inches(3.4), Inches(sy + 0.03), Inches(bar_w), Inches(0.18))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(C["purple"])
        bar.line.fill.background()

        pv = slide.shapes.add_textbox(Inches(3.4 + bar_w + 0.05), Inches(sy), Inches(1.5), Inches(0.25))
        tf_p = pv.text_frame
        p_p = tf_p.paragraphs[0]
        r_p = p_p.add_run()
        r_p.text = f"{fmt_currency(store['totalRevenue'])}  ({pct:.1f}%)"
        r_p.font.size = Pt(7)
        r_p.font.bold = True
        r_p.font.color.rgb = rgb(C["purple"])
        r_p.font.name = "Calibri"


def build_ranking_table(prs, stores, month_year, total_slides, start_page):
    TABLE_TOP = 1.55
    FOOTER_Y = 5.33
    ROW_H = 0.28
    ROWS_PER_PAGE = int(math.floor((FOOTER_Y - TABLE_TOP - 0.3) / ROW_H))

    total_rev = sum(s["totalRevenue"] for s in stores)
    chunks = [stores[i:i+ROWS_PER_PAGE] for i in range(0, len(stores), ROWS_PER_PAGE)]
    total_pages = len(chunks)
    pages_built = 0

    headers = ["Rank", "Store Name", "Revenue", "Oil Changes", "Avg Rev/Inv", "Share %"]

    for page_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_background(slide)
        subtitle = f"All locations ranked by total revenue — {month_year}"
        if total_pages > 1:
            subtitle += f"  ({page_idx+1} of {total_pages})"
        add_slide_header(slide, "Store Performance Ranking", subtitle)
        add_footer(slide, start_page + page_idx, total_slides)

        rows = len(chunk) + 1
        cols = len(headers)
        tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(TABLE_TOP), Inches(9.0), Inches(rows * ROW_H)).table
        col_widths = [Inches(0.6), Inches(3.0), Inches(1.4), Inches(1.2), Inches(1.4), Inches(1.0)]
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = cw

        for ci, ht in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = ht
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = rgb(C["white"])
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.LEFT if ci < 2 else PP_ALIGN.RIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(C["purple"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for ri, store in enumerate(chunk):
            pct = store["totalRevenue"] / total_rev * 100 if total_rev else 0
            row_data = [
                f"#{store['rank']}",
                store["name"],
                fmt_currency(store["totalRevenue"]),
                fmt_number(store["invoices"]),
                f"${store['avgRevPerInvoice']:.2f}",
                f"{pct:.1f}%",
            ]
            bg = C["white"] if ri % 2 == 0 else C["offWhite"]
            for ci, val in enumerate(row_data):
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = "Calibri"
                    paragraph.alignment = PP_ALIGN.LEFT if ci < 2 else PP_ALIGN.RIGHT
                    if ci == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = rgb(C["gold"])
                    else:
                        paragraph.font.color.rgb = rgb(C["darkGray"])
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(bg)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        pages_built += 1

    return pages_built


def build_performance_matrix(prs, stores, month_year, total_slides, start_page):
    TABLE_TOP = 1.55
    FOOTER_Y = 5.33
    ROW_H = 0.28
    ROWS_PER_PAGE = int(math.floor((FOOTER_Y - TABLE_TOP - 0.3) / ROW_H))

    chunks = [stores[i:i+ROWS_PER_PAGE] for i in range(0, len(stores), ROWS_PER_PAGE)]
    total_pages = len(chunks)
    pages_built = 0

    headers = ["Store", "Invoices", "Avg Rev/Inv", "Revenue", "Vehicles"]

    for page_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_background(slide)
        subtitle = f"Volume vs. average ticket — {month_year}"
        if total_pages > 1:
            subtitle += f"  ({page_idx+1} of {total_pages})"
        add_slide_header(slide, "Store Performance Matrix", subtitle)
        add_footer(slide, start_page + page_idx, total_slides)

        rows = len(chunk) + 1
        cols = len(headers)
        tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(TABLE_TOP), Inches(9.0), Inches(rows * ROW_H)).table
        col_widths = [Inches(3.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = cw

        for ci, ht in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = ht
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = rgb(C["white"])
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(C["purple"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for ri, store in enumerate(chunk):
            row_data = [
                store["name"],
                fmt_number(store["invoices"]),
                f"${store['avgRevPerInvoice']:.2f}",
                fmt_currency(store["totalRevenue"]),
                fmt_number(store["vehicles"]),
            ]
            bg = C["white"] if ri % 2 == 0 else C["offWhite"]
            for ci, val in enumerate(row_data):
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = "Calibri"
                    paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
                    paragraph.font.color.rgb = rgb(C["darkGray"])
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(bg)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        pages_built += 1

    return pages_built


def build_product_mix(prs, stores, month_year, total_slides, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_slide_header(slide, "Product Mix Analysis", f"Revenue by product category — {month_year}")
    add_footer(slide, page_num, total_slides)

    total_rev = sum(s["totalRevenue"] for s in stores)
    cat_rev = {}
    for s in stores:
        for pb in s["productBreakdown"]:
            cat = pb["category"]
            cat_rev[cat] = cat_rev.get(cat, 0) + pb["revenue"]

    rs_rev = sum(v for k, v in cat_rev.items() if k in ("Royal Purple High", "Royal Purple Syn"))
    hmx_rev = sum(v for k, v in cat_rev.items() if k in ("High Mileage", "High Mileage Syn"))
    other_rev = total_rev - rs_rev - hmx_rev
    top3 = [
        ("Royal Purple Synthetic", rs_rev),
        ("High Mileage", hmx_rev),
        ("Other / Specialty", max(other_rev, 0)),
    ]

    for i, (cat, rev) in enumerate(top3):
        x = 0.5 + i * 3.1
        pct = rev / total_rev * 100 if total_rev else 0
        _add_metric(slide, x, 1.55, 2.8, f"{pct:.1f}%", cat, fmt_currency(rev))

    _add_thin_divider(slide, 0.5, 2.55, 9.0)

    sku_rev = {}
    for s in stores:
        for pb in s["productBreakdown"]:
            sku_rev[pb["code"]] = sku_rev.get(pb["code"], 0) + pb["revenue"]

    sorted_skus = sorted(sku_rev.items(), key=lambda x: -x[1])[:10]
    max_sku_rev = sorted_skus[0][1] if sorted_skus else 1

    lbl_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4), Inches(0.3))
    tf = lbl_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Top 10 Product SKUs"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = rgb(C["purple"])
    r.font.name = "Calibri"

    for i, (code, rev) in enumerate(sorted_skus):
        y = 3.05 + i * 0.23
        bar_w = max(0.2, (rev / max_sku_rev) * 4.5)

        nm = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.8), Inches(0.2))
        tf_n = nm.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.RIGHT
        r_n = p_n.add_run()
        r_n.text = get_product_display_name(code)
        r_n.font.size = Pt(7)
        r_n.font.color.rgb = rgb(C["darkGray"])
        r_n.font.name = "Calibri"

        bar = slide.shapes.add_shape(1, Inches(3.4), Inches(y + 0.03), Inches(bar_w), Inches(0.14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(C["purple"])
        bar.line.fill.background()

        vb = slide.shapes.add_textbox(Inches(3.4 + bar_w + 0.05), Inches(y), Inches(1.5), Inches(0.2))
        tf_v = vb.text_frame
        p_v = tf_v.paragraphs[0]
        r_v = p_v.add_run()
        r_v.text = fmt_currency(rev)
        r_v.font.size = Pt(7)
        r_v.font.bold = True
        r_v.font.color.rgb = rgb(C["purple"])
        r_v.font.name = "Calibri"


def _aggregate_product_categories(stores):
    cats = {}
    for s in stores:
        for pb in s["productBreakdown"]:
            cat = pb["category"]
            if cat not in cats:
                cats[cat] = {
                    "category": cat,
                    "totalRevenue": 0,
                    "totalLineCount": 0,
                    "storeCount": 0,
                    "stores": {},
                    "skus": {},
                }
            cats[cat]["totalRevenue"] += pb["revenue"]
            cats[cat]["totalLineCount"] += pb.get("lineCount", 0)

            sname = s["name"]
            if sname not in cats[cat]["stores"]:
                cats[cat]["stores"][sname] = {"revenue": 0, "lineCount": 0}
                cats[cat]["storeCount"] += 1
            cats[cat]["stores"][sname]["revenue"] += pb["revenue"]
            cats[cat]["stores"][sname]["lineCount"] += pb.get("lineCount", 0)

            code = pb["code"]
            if code not in cats[cat]["skus"]:
                cats[cat]["skus"][code] = {"revenue": 0, "lineCount": 0}
            cats[cat]["skus"][code]["revenue"] += pb["revenue"]
            cats[cat]["skus"][code]["lineCount"] += pb.get("lineCount", 0)

    result = sorted(cats.values(), key=lambda c: -c["totalRevenue"])
    return result


def build_product_deep_dive(prs, cat_data, stores, month_year, total_slides, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)

    cat_name = cat_data["category"]
    total_net_rev = sum(s["totalRevenue"] for s in stores)
    cat_pct = cat_data["totalRevenue"] / total_net_rev * 100 if total_net_rev else 0
    total_inv = sum(s["invoices"] for s in stores)

    add_slide_header(
        slide,
        f"Product Deep Dive: {cat_name}",
        f"{cat_pct:.1f}% of network revenue — {month_year}",
    )
    add_footer(slide, page_num, total_slides)

    is_mc = cat_name == "Max-Clean"

    metrics = [
        (fmt_currency(cat_data["totalRevenue"]), "Category Revenue", f"{cat_pct:.1f}% of network"),
        (fmt_number(cat_data["totalLineCount"]), "Transactions", "Product line items"),
        (fmt_number(cat_data["storeCount"]), "Stores Selling", f"of {len(stores)} locations"),
    ]

    if is_mc:
        mc_total = sum(s.get("maxClean", {}).get("total", 0) for s in stores)
        mc_rate = mc_total / total_inv * 100 if total_inv else 0
        metrics.append((f"{mc_rate:.1f}%", "Attachment Rate", f"{fmt_number(mc_total)} of {fmt_number(total_inv)} inv"))
    else:
        avg_per_store = cat_data["totalRevenue"] / cat_data["storeCount"] if cat_data["storeCount"] else 0
        metrics.append((fmt_currency(avg_per_store), "Avg per Store", "Revenue per location"))

    mw = 2.1
    gap = 0.15
    sx = 0.5
    for i, (val, lbl, sub) in enumerate(metrics):
        _add_metric(slide, sx + i * (mw + gap), 1.55, mw, val, lbl, sub)

    _add_thin_divider(slide, 0.5, 2.55, 9.0)

    sorted_stores = sorted(cat_data["stores"].items(), key=lambda x: -x[1]["revenue"])
    top_stores = sorted_stores[:8]

    if top_stores:
        max_store_rev = top_stores[0][1]["revenue"]

        lbl = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4), Inches(0.25))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{cat_name} Revenue by Store"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = rgb(C["purple"])
        r.font.name = "Calibri"

        for pi, (sname, sdata) in enumerate(top_stores):
            py = 3.05 + pi * 0.28
            bar_w = max(0.2, (sdata["revenue"] / max_store_rev) * 3.0) if max_store_rev else 0.2

            nm = slide.shapes.add_textbox(Inches(0.5), Inches(py), Inches(2.3), Inches(0.24))
            tf_n = nm.text_frame
            tf_n.word_wrap = False
            p_n = tf_n.paragraphs[0]
            p_n.alignment = PP_ALIGN.RIGHT
            r_n = p_n.add_run()
            display_name = sname if len(sname) <= 22 else sname[:20] + ".."
            r_n.text = display_name
            r_n.font.size = Pt(7)
            r_n.font.color.rgb = rgb(C["darkGray"])
            r_n.font.name = "Calibri"

            bar = slide.shapes.add_shape(1, Inches(2.9), Inches(py + 0.04), Inches(bar_w), Inches(0.16))
            bar.fill.solid()
            bar.fill.fore_color.rgb = rgb(C["purple"])
            bar.line.fill.background()

            store_pct = sdata["revenue"] / cat_data["totalRevenue"] * 100 if cat_data["totalRevenue"] else 0
            vb = slide.shapes.add_textbox(Inches(2.9 + bar_w + 0.05), Inches(py), Inches(1.5), Inches(0.24))
            tf_v = vb.text_frame
            p_v = tf_v.paragraphs[0]
            r_v = p_v.add_run()
            r_v.text = f"{fmt_currency(sdata['revenue'])} ({store_pct:.0f}%)"
            r_v.font.size = Pt(7)
            r_v.font.bold = True
            r_v.font.color.rgb = rgb(C["purple"])
            r_v.font.name = "Calibri"

    desc = PRODUCT_DESCRIPTIONS.get(cat_name, "")
    top_store_name = sorted_stores[0][0] if sorted_stores else "N/A"
    top_store_pct = (sorted_stores[0][1]["revenue"] / cat_data["totalRevenue"] * 100) if sorted_stores and cat_data["totalRevenue"] else 0

    if is_mc:
        mc_with_rp = sum(s.get("maxClean", {}).get("withRpOil", 0) for s in stores)
        mc_non_rp = sum(s.get("maxClean", {}).get("withNonRpOil", 0) for s in stores)
        mc_total_inv = sum(s.get("maxClean", {}).get("total", 0) for s in stores)
        best_rate_store = max(stores, key=lambda s: s.get("maxClean", {}).get("attachmentRate", 0))
        best_lift_store = max(stores, key=lambda s: s.get("maxClean", {}).get("ticketLift", 0))
        br_mc = best_rate_store.get("maxClean", {})
        bl_mc = best_lift_store.get("maxClean", {})

        rp_pct = mc_with_rp / mc_total_inv * 100 if mc_total_inv else 0
        note_text = (
            f"{desc}\n\n"
            f"Of {fmt_number(mc_total_inv)} Max-Clean transactions, "
            f"{fmt_number(mc_with_rp)} ({rp_pct:.0f}%) were paired with RP oil "
            f"and {fmt_number(mc_non_rp)} ({100 - rp_pct:.0f}%) with non-RP oil (upsell opportunity). "
            f"{best_rate_store['name']} leads attachment at {br_mc.get('attachmentRate', 0):.1f}%. "
            f"{best_lift_store['name']} delivers the highest ticket lift at +${bl_mc.get('ticketLift', 0):.2f}."
        )
    else:
        low_stores = [sn for sn, sd in sorted_stores if sd["revenue"] < cat_data["totalRevenue"] / max(len(sorted_stores), 1) * 0.5]
        opp_note = f" {len(low_stores)} store(s) sell below average — potential growth targets." if low_stores else ""
        note_text = (
            f"{desc}\n\n"
            f"{top_store_name} leads with {top_store_pct:.0f}% of {cat_name} revenue across the network. "
            f"This product is sold at {cat_data['storeCount']} of {len(stores)} locations."
            f"{opp_note}"
        )

    sorted_skus = sorted(cat_data["skus"].items(), key=lambda x: -x[1]["revenue"])[:5]
    if sorted_skus and len(sorted_skus) > 1:
        sku_lines = "\n\nTop SKUs: " + ", ".join(
            f"{get_product_display_name(code)} ({fmt_currency(sd['revenue'])})"
            for code, sd in sorted_skus
        )
        note_text += sku_lines

    n_lbl = slide.shapes.add_textbox(Inches(6.5), Inches(2.7), Inches(3.2), Inches(0.25))
    tf_nl = n_lbl.text_frame
    p_nl = tf_nl.paragraphs[0]
    r_nl = p_nl.add_run()
    r_nl.text = "Insights"
    r_nl.font.size = Pt(11)
    r_nl.font.bold = True
    r_nl.font.color.rgb = rgb(C["purple"])
    r_nl.font.name = "Calibri"

    n_body = slide.shapes.add_textbox(Inches(6.5), Inches(3.0), Inches(3.2), Inches(2.1))
    tf_nb = n_body.text_frame
    tf_nb.word_wrap = True
    p_nb = tf_nb.paragraphs[0]
    r_nb = p_nb.add_run()
    r_nb.text = note_text
    r_nb.font.size = Pt(8)
    r_nb.font.color.rgb = rgb(C["darkGray"])
    r_nb.font.name = "Calibri"


def build_section_divider(prs, title, subtitle, total_slides, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, C["dark"])

    if os.path.exists(BG_NEVER_SETTLE):
        slide.shapes.add_picture(
            BG_NEVER_SETTLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
        )
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb(C["dark"])
        _set_shape_alpha(overlay, 70000)
        overlay.line.fill.background()

    bar = slide.shapes.add_shape(1, Inches(2.5), Inches(1.8), Inches(5), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(C["gold"])
    bar.line.fill.background()

    t_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(0.7))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = rgb(C["white"])
    r.font.name = "Calibri"

    s_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(8), Inches(0.4))
    tf2 = s_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(12)
    r2.font.color.rgb = rgb(C["purpleLight"])
    r2.font.name = "Calibri"

    ns_box = slide.shapes.add_textbox(Inches(1), Inches(3.15), Inches(8), Inches(0.35))
    tf_ns = ns_box.text_frame
    p_ns = tf_ns.paragraphs[0]
    p_ns.alignment = PP_ALIGN.CENTER
    r_ns = p_ns.add_run()
    r_ns.text = "NEVER SETTLE"
    r_ns.font.size = Pt(14)
    r_ns.font.bold = True
    r_ns.font.color.rgb = rgb(C["goldLight"])
    r_ns.font.name = "Calibri"

    add_footer(slide, page_num, total_slides)


def build_deep_dive(prs, store, stores, month_year, total_slides, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_slide_header(
        slide,
        f"Store Deep Dive: {store['name']}",
        f"Rank #{store['rank']} of {len(stores)} by Revenue | {month_year}"
    )
    add_footer(slide, page_num, total_slides)

    metrics = [
        (fmt_currency(store["totalRevenue"]), "Total Revenue", month_year),
        (fmt_number(store["invoices"]), "Oil Changes", "Invoice count"),
        (f"${store['avgRevPerInvoice']:.2f}", "Avg Rev/Invoice", "Per transaction"),
        (f"#{store['rank']}", "Network Rank", f"of {len(stores)} stores"),
    ]
    mw = 2.1
    gap = 0.15
    sx = 0.5
    for i, (val, lbl, sub) in enumerate(metrics):
        _add_metric(slide, sx + i * (mw + gap), 1.55, mw, val, lbl, sub)

    _add_thin_divider(slide, 0.5, 2.55, 9.0)

    top_products = store["productBreakdown"][:6]
    if top_products:
        max_prod_rev = top_products[0]["revenue"]

        lbl = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4), Inches(0.25))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Revenue by Product"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = rgb(C["purple"])
        r.font.name = "Calibri"

        for pi, prod in enumerate(top_products):
            py = 3.05 + pi * 0.28
            bar_w = max(0.2, (prod["revenue"] / max_prod_rev) * 3.0) if max_prod_rev else 0.2

            nm = slide.shapes.add_textbox(Inches(0.5), Inches(py), Inches(2.3), Inches(0.24))
            tf_n = nm.text_frame
            tf_n.word_wrap = False
            p_n = tf_n.paragraphs[0]
            p_n.alignment = PP_ALIGN.RIGHT
            r_n = p_n.add_run()
            r_n.text = get_product_display_name(prod["code"])
            r_n.font.size = Pt(7)
            r_n.font.color.rgb = rgb(C["darkGray"])
            r_n.font.name = "Calibri"

            bar = slide.shapes.add_shape(1, Inches(2.9), Inches(py + 0.04), Inches(bar_w), Inches(0.16))
            bar.fill.solid()
            bar.fill.fore_color.rgb = rgb(C["purple"])
            bar.line.fill.background()

            vb = slide.shapes.add_textbox(Inches(2.9 + bar_w + 0.05), Inches(py), Inches(1.5), Inches(0.24))
            tf_v = vb.text_frame
            p_v = tf_v.paragraphs[0]
            r_v = p_v.add_run()
            r_v.text = fmt_currency(prod["revenue"])
            r_v.font.size = Pt(7)
            r_v.font.bold = True
            r_v.font.color.rgb = rgb(C["purple"])
            r_v.font.name = "Calibri"

    total_rev = sum(s["totalRevenue"] for s in stores)
    pct = store["totalRevenue"] / total_rev * 100 if total_rev else 0
    cat_breakdown = {}
    for pb in store["productBreakdown"]:
        cat_breakdown[pb["category"]] = cat_breakdown.get(pb["category"], 0) + pb["revenue"]
    top_cat = max(cat_breakdown.items(), key=lambda x: x[1])[0] if cat_breakdown else "N/A"
    top_cat_pct = (cat_breakdown.get(top_cat, 0) / store["totalRevenue"] * 100) if store["totalRevenue"] else 0

    top_cat_desc = get_product_category_desc(top_cat)
    desc_line = f" {top_cat_desc}" if top_cat_desc else ""

    note_text = (
        f"This location contributes {pct:.1f}% of total network revenue. "
        f"With {fmt_number(store['invoices'])} oil changes and an average ticket of "
        f"${store['avgRevPerInvoice']:.2f}, {store['name']} "
        f"{'leads' if store['rank'] == 1 else 'ranks #' + str(store['rank'])} in the network. "
        f"{top_cat} products represent {top_cat_pct:.0f}% of this location's mix."
        f"{desc_line}\n\nTop Product: {store['topProduct']}"
    )

    n_lbl = slide.shapes.add_textbox(Inches(6.5), Inches(2.7), Inches(3.2), Inches(0.25))
    tf_nl = n_lbl.text_frame
    p_nl = tf_nl.paragraphs[0]
    r_nl = p_nl.add_run()
    r_nl.text = "Location Notes"
    r_nl.font.size = Pt(11)
    r_nl.font.bold = True
    r_nl.font.color.rgb = rgb(C["purple"])
    r_nl.font.name = "Calibri"

    n_body = slide.shapes.add_textbox(Inches(6.5), Inches(3.0), Inches(3.2), Inches(2.1))
    tf_nb = n_body.text_frame
    tf_nb.word_wrap = True
    p_nb = tf_nb.paragraphs[0]
    r_nb = p_nb.add_run()
    r_nb.text = note_text
    r_nb.font.size = Pt(8)
    r_nb.font.color.rgb = rgb(C["darkGray"])
    r_nb.font.name = "Calibri"


def build_next_steps(prs, stores, month_year, total_slides, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_slide_header(slide, "Next Steps & Recommendations", f"Action items for {month_year}")
    add_footer(slide, page_num, total_slides)

    bottom = stores[-1] if stores else None
    top = stores[0] if stores else None
    actions = [
        (
            "Boost Underperformers",
            f"Target {bottom['name']} with promotional offers to increase volume and average ticket size." if bottom else "Identify underperforming locations for targeted promotions.",
        ),
        (
            "Expand High Mileage",
            "High Mileage products show strong margins — push upsell training to increase penetration across all locations.",
        ),
        (
            "Replicate Top Performer",
            f"Study {top['name']}'s practices and replicate successful strategies across the network." if top else "Analyze top store practices for network-wide adoption.",
        ),
        (
            "Monthly Trend Tracking",
            "Establish month-over-month tracking to identify seasonal patterns and measure the impact of promotional campaigns.",
        ),
    ]

    for i, (title, body) in enumerate(actions):
        y = 1.55 + i * 0.92
        _add_thin_divider(slide, 0.5, y, 9.0)

        nb = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.08), Inches(0.35), Inches(0.3))
        tf_n = nb.text_frame
        p_n = tf_n.paragraphs[0]
        r_n = p_n.add_run()
        r_n.text = str(i + 1)
        r_n.font.size = Pt(16)
        r_n.font.bold = True
        r_n.font.color.rgb = rgb(C["gold"])
        r_n.font.name = "Calibri"

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.08), Inches(8.5), Inches(0.25))
        tf_t = tb.text_frame
        p_t = tf_t.paragraphs[0]
        r_t = p_t.add_run()
        r_t.text = title
        r_t.font.size = Pt(12)
        r_t.font.bold = True
        r_t.font.color.rgb = rgb(C["purple"])
        r_t.font.name = "Calibri"

        bb = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.36), Inches(8.5), Inches(0.5))
        tf_b = bb.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        r_b = p_b.add_run()
        r_b.text = body
        r_b.font.size = Pt(9)
        r_b.font.color.rgb = rgb(C["darkGray"])
        r_b.font.name = "Calibri"


def build_closing_slide(prs, stores, month_year, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, C["dark"])

    if os.path.exists(BG_NEVER_SETTLE):
        slide.shapes.add_picture(
            BG_NEVER_SETTLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
        )
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb(C["dark"])
        _set_shape_alpha(overlay, 60000)
        overlay.line.fill.background()

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(C["gold"])
    bar.line.fill.background()

    s_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.4))
    tf2 = s_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"Partnership Hub Report | {month_year}"
    r2.font.size = Pt(12)
    r2.font.color.rgb = rgb(C["purpleLight"])
    r2.font.name = "Calibri"

    total_rev = sum(s["totalRevenue"] for s in stores)
    total_inv = sum(s["invoices"] for s in stores)
    avg_rev = total_rev / total_inv if total_inv else 0

    stat_x = 1.5
    stat_y = 2.2
    stat_w = 7.0
    stat_h = 2.0

    stat_bg = slide.shapes.add_shape(1, Inches(stat_x), Inches(stat_y), Inches(stat_w), Inches(stat_h))
    stat_bg.fill.solid()
    stat_bg.fill.fore_color.rgb = rgb(C["purple"])
    stat_bg.line.fill.background()

    gold_top = slide.shapes.add_shape(1, Inches(stat_x), Inches(stat_y), Inches(stat_w), Inches(0.06))
    gold_top.fill.solid()
    gold_top.fill.fore_color.rgb = rgb(C["gold"])
    gold_top.line.fill.background()

    lbl = slide.shapes.add_textbox(Inches(stat_x + 0.3), Inches(stat_y + 0.15), Inches(stat_w - 0.6), Inches(0.25))
    tf_l = lbl.text_frame
    p_l = tf_l.paragraphs[0]
    p_l.alignment = PP_ALIGN.CENTER
    r_l = p_l.add_run()
    r_l.text = f"{month_year} Summary"
    r_l.font.size = Pt(11)
    r_l.font.bold = True
    r_l.font.color.rgb = rgb(C["goldLight"])
    r_l.font.name = "Calibri"

    items = [
        (fmt_currency(total_rev), "Total Revenue"),
        (fmt_number(total_inv), "Oil Changes"),
        (f"${avg_rev:.2f}", "Avg Rev/Invoice"),
        (stores[0]["name"] if stores else "N/A", "Top Store"),
    ]
    col_w = stat_w / 4
    for i, (val, lab) in enumerate(items):
        ix = stat_x + i * col_w + 0.15
        iy = stat_y + 0.55

        vb = slide.shapes.add_textbox(Inches(ix), Inches(iy), Inches(col_w - 0.3), Inches(0.4))
        tf_v = vb.text_frame
        p_v = tf_v.paragraphs[0]
        p_v.alignment = PP_ALIGN.CENTER
        r_v = p_v.add_run()
        r_v.text = val
        font_sz = 16 if len(val) > 10 else 20
        r_v.font.size = Pt(font_sz)
        r_v.font.bold = True
        r_v.font.color.rgb = rgb(C["white"])
        r_v.font.name = "Calibri"

        lb = slide.shapes.add_textbox(Inches(ix), Inches(iy + 0.45), Inches(col_w - 0.3), Inches(0.2))
        tf_lb = lb.text_frame
        p_lb = tf_lb.paragraphs[0]
        p_lb.alignment = PP_ALIGN.CENTER
        r_lb = p_lb.add_run()
        r_lb.text = lab
        r_lb.font.size = Pt(8)
        r_lb.font.color.rgb = rgb(C["purpleLight"])
        r_lb.font.name = "Calibri"

    ns_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.35))
    tf_ns = ns_box.text_frame
    p_ns = tf_ns.paragraphs[0]
    p_ns.alignment = PP_ALIGN.CENTER
    r_ns = p_ns.add_run()
    r_ns.text = "NEVER SETTLE"
    r_ns.font.size = Pt(16)
    r_ns.font.bold = True
    r_ns.font.color.rgb = rgb(C["goldLight"])
    r_ns.font.name = "Calibri"

    from datetime import datetime as _dt
    prepared_date = _dt.now().strftime("%B %d, %Y")

    contact = slide.shapes.add_textbox(Inches(1), Inches(4.65), Inches(8), Inches(0.25))
    tf_c = contact.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    r_c = p_c.add_run()
    r_c.text = f"Prepared {prepared_date}"
    r_c.font.size = Pt(9)
    r_c.font.color.rgb = rgb(C["purpleLight"])
    r_c.font.name = "Calibri"

    contact2 = slide.shapes.add_textbox(Inches(1), Inches(4.9), Inches(8), Inches(0.25))
    tf_c2 = contact2.text_frame
    p_c2 = tf_c2.paragraphs[0]
    p_c2.alignment = PP_ALIGN.CENTER
    r_c2 = p_c2.add_run()
    r_c2.text = "ThrottlePro — More Cars. More Loyalty. Less Stress."
    r_c2.font.size = Pt(8)
    r_c2.font.color.rgb = rgb(C["midGray"])
    r_c2.font.name = "Calibri"

    add_footer(slide, total_slides, total_slides)


def build_distribution_map_slide(prs, map_image_path, title, total_slides, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, C["offWhite"])

    hdr_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    hdr_bar.fill.solid()
    hdr_bar.fill.fore_color.rgb = rgb(C["purple"])
    hdr_bar.line.fill.background()

    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7), Inches(0.45))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = rgb(C["white"])
    r.font.name = "Calibri"

    s_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(7), Inches(0.3))
    tf2 = s_box.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Territory Map"
    r2.font.size = Pt(11)
    r2.font.color.rgb = rgb(C["goldLight"])
    r2.font.name = "Calibri"

    add_royal_purple_badge(slide)

    map_top = 1.2
    map_bottom = 5.2
    available_h = map_bottom - map_top
    available_w = 9.0
    margin_x = 0.5

    from PIL import Image as PILImage
    try:
        img = PILImage.open(map_image_path)
        img_w, img_h = img.size
        aspect = img_w / img_h
        fit_w = available_w
        fit_h = fit_w / aspect
        if fit_h > available_h:
            fit_h = available_h
            fit_w = fit_h * aspect
        center_x = margin_x + (available_w - fit_w) / 2
        center_y = map_top + (available_h - fit_h) / 2
        slide.shapes.add_picture(
            map_image_path, Inches(center_x), Inches(center_y), Inches(fit_w), Inches(fit_h)
        )
    except Exception:
        slide.shapes.add_picture(
            map_image_path, Inches(margin_x), Inches(map_top), Inches(available_w), Inches(available_h)
        )

    add_footer(slide, page_num, total_slides)


def calculate_total_slides(num_stores, num_maps=0, num_product_cats=0):
    TABLE_TOP = 1.55
    FOOTER_Y = 5.33
    ROW_H = 0.28
    ROWS_PER_PAGE_RANK = int(math.floor((FOOTER_Y - TABLE_TOP - 0.3) / ROW_H))

    ROWS_PER_PAGE_MATRIX = ROWS_PER_PAGE_RANK

    rank_pages = math.ceil(num_stores / ROWS_PER_PAGE_RANK) if num_stores > 0 else 1
    matrix_pages = math.ceil(num_stores / ROWS_PER_PAGE_MATRIX) if num_stores > 0 else 1

    product_dive_slides = (1 + num_product_cats) if num_product_cats > 0 else 0

    total = (
        1 +  # cover
        1 +  # toc
        1 +  # exec summary kpis
        1 +  # exec observations
        1 +  # revenue overview
        rank_pages +  # ranking table
        matrix_pages +  # performance matrix
        1 +  # product mix
        product_dive_slides +  # product deep dives (divider + per-category)
        num_maps +  # map image slides
        1 +  # section divider (store deep dives)
        num_stores +  # store deep dives
        1 +  # next steps
        1    # closing
    )
    return total


def generate_report(file_path, output_path=None, map_images=None):
    stores, month_year = parse_excel(file_path)
    maps = map_images or []

    product_cats = _aggregate_product_categories(stores)

    if not output_path:
        month_abbr = month_year.split()[0][:3] if month_year else "Jan"
        year = month_year.split()[-1] if month_year else "2025"
        output_path = f"Royal_Purple_Partnership_Report_{month_abbr}{year}.pptx"

    total_slides = calculate_total_slides(len(stores), len(maps), len(product_cats))

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    build_cover_slide(prs, stores, month_year, total_slides)
    build_toc_slide(prs, total_slides)
    build_exec_summary_kpis(prs, stores, month_year, total_slides)
    build_exec_observations(prs, stores, month_year, total_slides)
    build_revenue_overview(prs, stores, month_year, total_slides)

    page = 6
    rank_pages = build_ranking_table(prs, stores, month_year, total_slides, page)
    page += rank_pages

    matrix_pages = build_performance_matrix(prs, stores, month_year, total_slides, page)
    page += matrix_pages

    build_product_mix(prs, stores, month_year, total_slides, page)
    page += 1

    if product_cats:
        build_section_divider(prs, "Product Deep Dives", f"Category-level performance analysis — {month_year}", total_slides, page)
        page += 1

        for cat_data in product_cats:
            build_product_deep_dive(prs, cat_data, stores, month_year, total_slides, page)
            page += 1

    for i, map_info in enumerate(maps):
        map_path = map_info.get("path", map_info) if isinstance(map_info, dict) else map_info
        map_title = map_info.get("title", f"Map {i + 1}") if isinstance(map_info, dict) else f"Map {i + 1}"
        build_distribution_map_slide(prs, map_path, map_title, total_slides, page)
        page += 1

    build_section_divider(prs, "Store-Level Deep Dives", f"Individual performance analysis — {month_year}", total_slides, page)
    page += 1

    for store in stores:
        build_deep_dive(prs, store, stores, month_year, total_slides, page)
        page += 1

    build_next_steps(prs, stores, month_year, total_slides, page)
    page += 1

    build_closing_slide(prs, stores, month_year, total_slides)

    prs.save(output_path)
    return output_path, stores, month_year
